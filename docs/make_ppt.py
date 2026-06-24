"""
WBS 管理系統投影片
主題：白色簡約，敘事線：單頁原型 → 解決痛點 → 邁向正式產品
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SS   = os.path.join(BASE, "screenshots")
OUT  = os.path.join(BASE, "presentation.pptx")

W = Inches(13.33)
H = Inches(7.5)

# ── 色票（白色簡約） ──────────────────────────────────────
BG      = RGBColor(0xFF, 0xFF, 0xFF)   # 白底
BG2     = RGBColor(0xF8, 0xFA, 0xFC)  # 淺灰底
CARD    = RGBColor(0xF1, 0xF5, 0xF9)  # 卡片背景
BORDER  = RGBColor(0xE2, 0xE8, 0xF0)  # 分隔線
BLUE    = RGBColor(0x25, 0x63, 0xEB)  # 主藍
LBLUE   = RGBColor(0xDB, 0xEA, 0xFE)  # 淺藍
DARK    = RGBColor(0x0F, 0x17, 0x2A)  # 深色標題
BODY    = RGBColor(0x33, 0x41, 0x55)  # 內文
MUTED   = RGBColor(0x64, 0x74, 0x8B)  # 次要文字
GREEN   = RGBColor(0x16, 0xA3, 0x4A)  # 成功綠
LGREEN  = RGBColor(0xDC, 0xFC, 0xE7)  # 淺綠
RED     = RGBColor(0xDC, 0x26, 0x26)  # 錯誤紅
LRED    = RGBColor(0xFE, 0xE2, 0xE2)  # 淺紅
AMBER   = RGBColor(0xD9, 0x77, 0x06)  # 警告橙
LAMBER  = RGBColor(0xFE, 0xF3, 0xC7)  # 淺橙
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

TOTAL = 13

# ── 工具函式 ──────────────────────────────────────────────

def rect(sl, l, t, w, h, fill=BG, line_col=None, lw=Pt(0.75)):
    s = sl.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_col:
        s.line.color.rgb = line_col; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(sl, text, l, t, w, h,
        size=Pt(14), bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False):
    txb = sl.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def img(sl, fname, l, t, w=None, h=None):
    path = os.path.join(SS, fname)
    if not os.path.exists(path): return
    kw = {}
    if w: kw["width"]  = w
    if h: kw["height"] = h
    sl.shapes.add_picture(path, l, t, **kw)

def new_slide(bg=BG):
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, W, H, fill=bg)
    return sl

def header(sl, title, subtitle=""):
    """白底頂欄，藍色左側色條"""
    rect(sl, 0, 0, W, Inches(1.05), fill=BG)
    rect(sl, 0, Inches(1.05), W, Pt(1.5), fill=BORDER)
    rect(sl, 0, 0, Inches(0.22), Inches(1.05), fill=BLUE)
    txt(sl, title,
        Inches(0.38), Inches(0.12), Inches(10), Inches(0.58),
        size=Pt(24), bold=True, color=DARK)
    if subtitle:
        txt(sl, subtitle,
            Inches(0.38), Inches(0.65), Inches(12), Inches(0.34),
            size=Pt(12), color=MUTED)

def pnum(sl, n):
    txt(sl, f"{n} / {TOTAL}",
        Inches(12.5), Inches(7.1), Inches(0.8), Inches(0.35),
        size=Pt(10), color=MUTED, align=PP_ALIGN.RIGHT)

def tag(sl, label, l, t, w, h, bg=LBLUE, fg=BLUE):
    rect(sl, l, t, w, h, fill=bg)
    txt(sl, label, l, t, w, h,
        size=Pt(11), bold=True, color=fg, align=PP_ALIGN.CENTER)

def check(val):
    if val == "✅":   return GREEN, LGREEN
    if val == "❌":   return RED,   LRED
    if val == "⚠️":  return AMBER,  LAMBER
    return DARK, CARD

# ═══════════════════════════════════════════════════════════
# 1. 封面
# ═══════════════════════════════════════════════════════════
sl = new_slide(BG2)
rect(sl, 0, 0, Inches(0.45), H, fill=BLUE)
rect(sl, Inches(0.45), 0, W - Inches(0.45), H, fill=BG)

txt(sl, "WBS 管理系統",
    Inches(0.75), Inches(1.6), Inches(11.5), Inches(1.2),
    size=Pt(52), bold=True, color=DARK)
rect(sl, Inches(0.75), Inches(2.85), Inches(3.2), Pt(3), fill=BLUE)
txt(sl, "從單頁原型，走向企業協作平台",
    Inches(0.75), Inches(2.95), Inches(11), Inches(0.55),
    size=Pt(20), color=BODY)
txt(sl, "IT 部門 WBS 多人即時協作系統  ·  角色權限  ·  標準模板  ·  WebSocket 即時同步",
    Inches(0.75), Inches(3.65), Inches(11), Inches(0.4),
    size=Pt(13), color=MUTED)

tags_data = ["Java 21", "Spring Boot 3.4", "Vue 3", "WebSocket", "PostgreSQL 16"]
x = Inches(0.75)
for t_label in tags_data:
    tw = Inches(len(t_label) * 0.165 + 0.55)
    tag(sl, t_label, x, Inches(4.4), tw, Inches(0.38))
    x += tw + Inches(0.18)

txt(sl, "github.com/LinGaOscar/Spring-WbsScaff",
    Inches(0.75), Inches(6.85), Inches(6), Inches(0.38),
    size=Pt(11), color=MUTED, italic=True)
pnum(sl, 1)

# ═══════════════════════════════════════════════════════════
# 2. 起點：一切從一個需求開始
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "一切從一個需求開始", "最初只是想解決 IT 部門每次上線前的混亂")
pnum(sl, 2)

# 左側故事
story = [
    ("🗂  原點",
     "IT 部門每次新專案上線，都要重新整理一份 WBS，\n格式不一、版本混亂，找不到人負責哪個項目。"),
    ("💬  第一步",
     "從一個單頁面工具開始，只是想讓 WBS 可以快速\n編排，不再依賴 Excel 手工排版。"),
    ("🚀  演進",
     "隨著使用需求增加，逐步加入即時協作、角色權限、\n標準模板，從原型演進為可正式使用的平台。"),
]
for i, (title, desc) in enumerate(story):
    y = Inches(1.3 + i * 1.8)
    rect(sl, Inches(0.4), y, Inches(0.08), Inches(1.5), fill=BLUE)
    txt(sl, title,
        Inches(0.65), y, Inches(5.8), Inches(0.48),
        size=Pt(15), bold=True, color=DARK)
    txt(sl, desc,
        Inches(0.65), y + Inches(0.5), Inches(5.8), Inches(1.0),
        size=Pt(12), color=BODY)

# 右側：專案列表截圖
rect(sl, Inches(7.0), Inches(1.2), Inches(6.1), Inches(5.9),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-project-list.png", Inches(7.05), Inches(1.25), w=Inches(6.0))
txt(sl, "↑  專案列表頁面",
    Inches(7.0), Inches(6.95), Inches(6.1), Inches(0.38),
    size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 3. 痛點
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "這些問題，IT 團隊都遇過", "每次專案上線前的混亂，不應該一直重演")
pnum(sl, 3)

pains = [
    ("📋", "Excel 版本地獄",
     "N 個同事，N 個版本，沒有人知道哪份才是最新的"),
    ("🔁", "反覆開會確認進度",
     "主管無法即時掌握狀態，只能靠開會問人，效率極低"),
    ("👤", "責任歸屬不清",
     "任務沒有明確負責人與截止日，事情就這樣掉了"),
    ("📂", "文件格式不統一",
     "SIT / UAT / PROD 各階段文件每次都要重新想格式"),
    ("🔒", "資料缺乏隔離",
     "共用試算表沒有權限控管，跨科資料外洩風險高"),
]
for i, (icon, title, desc) in enumerate(pains):
    row, col = divmod(i, 3)
    x = Inches(0.38 + col * 4.35)
    y = Inches(1.25 + row * 2.6)
    rect(sl, x, y, Inches(4.15), Inches(2.3),
         fill=LRED, line_col=RGBColor(0xFC, 0xCA, 0xCA))
    rect(sl, x, y, Inches(0.07), Inches(2.3), fill=RED)
    txt(sl, f"{icon}  {title}",
        x + Inches(0.18), y + Inches(0.2),
        Inches(3.85), Inches(0.48),
        size=Pt(14), bold=True, color=RED)
    txt(sl, desc,
        x + Inches(0.18), y + Inches(0.75),
        Inches(3.85), Inches(1.35),
        size=Pt(12), color=BODY)

# ═══════════════════════════════════════════════════════════
# 4. 解方：三個核心能力
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "WBS 管理系統  —  三個核心能力")
pnum(sl, 4)

cores = [
    (BLUE, "01", "標準化模板",
     ["4 套系統模板（功能新增、EOS 下線、專案開發、系統開發）",
      "SIT → UAT → PROD 三階段架構預設到位",
      "30 秒建好完整 WBS，不需手工排版",
      "科別可另存自訂模板，保留最佳實踐"]),
    (GREEN, "02", "即時多人協作",
     ["WebSocket / STOMP 廣播，所有人看同一份資料",
      "延遲儲存機制，批次送出避免衝突",
      "協作游標顯示誰在哪個節點",
      "不需重新整理頁面即可看到他人異動"]),
    (RGBColor(0x7C, 0x3A, 0xED), "03", "精細角色權限",
     ["DIRECTOR / SECTION_CHIEF / PROJECT_LEADER / MEMBER 四層",
      "部門隔離：科與科之間資料互不可見",
      "歸檔保護：完成後全員唯讀，不可誤改",
      "部長可唯讀查閱所有下屬科進度"]),
]
for i, (col, num, title, bullets) in enumerate(cores):
    x = Inches(0.38 + i * 4.35)
    rect(sl, x, Inches(1.2), Inches(4.15), Inches(5.95), fill=CARD, line_col=BORDER)
    rect(sl, x, Inches(1.2), Inches(4.15), Inches(0.08), fill=col)
    lbl_w = Inches(0.55)
    rect(sl, x + Inches(0.18), Inches(1.38), lbl_w, lbl_w, fill=col)
    txt(sl, num,
        x + Inches(0.18), Inches(1.38), lbl_w, lbl_w,
        size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, title,
        x + Inches(0.9), Inches(1.44), Inches(3.1), Inches(0.52),
        size=Pt(16), bold=True, color=DARK)
    rect(sl, x + Inches(0.18), Inches(2.1), Inches(3.8), Pt(1), fill=BORDER)
    for j, b in enumerate(bullets):
        txt(sl, f"·  {b}",
            x + Inches(0.18), Inches(2.22 + j * 0.78),
            Inches(3.8), Inches(0.72),
            size=Pt(12), color=BODY)

# ═══════════════════════════════════════════════════════════
# 5. 功能展示：WBS 編輯器
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "實際畫面  —  WBS 編輯器", "功能新增模板套用，33 個節點一次到位")
pnum(sl, 5)

img(sl, "ppt-wbs-editor.png", Inches(3.8), Inches(1.15), w=Inches(9.3))
rect(sl, Inches(3.8), Inches(1.15), Inches(9.3), Inches(5.55),
     fill=None, line_col=BORDER)

annotations = [
    ("✅  SIT / UAT / PROD\n三階段 L1 結構",   BLUE,  Inches(1.7)),
    ("✅  L2 工作包\n+ L3 任務（含備註）",       GREEN, Inches(3.1)),
    ("✅  48 項快速子項\n左側拖曳面板",           RGBColor(0x7C,0x3A,0xED), Inches(4.5)),
    ("✅  狀態 / 負責人\n/ 日期欄位",             AMBER, Inches(5.9)),
]
for note, col, y in annotations:
    rect(sl, Inches(0.25), y, Inches(3.3), Inches(0.95),
         fill=CARD, line_col=col, lw=Pt(1.2))
    rect(sl, Inches(0.25), y, Inches(0.07), Inches(0.95), fill=col)
    txt(sl, note, Inches(0.4), y + Inches(0.08),
        Inches(3.0), Inches(0.8),
        size=Pt(11), bold=False, color=DARK)

# ═══════════════════════════════════════════════════════════
# 6. 功能展示：模板 + 快速子項
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "模板系統與快速子項", "標準化工作清單，大幅縮短新專案建置時間")
pnum(sl, 6)

# 模板截圖（左）
rect(sl, Inches(0.38), Inches(1.2), Inches(6.2), Inches(4.2),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-templates.png", Inches(0.43), Inches(1.25), w=Inches(6.1))
txt(sl, "主模板管理（4 套系統模板）",
    Inches(0.38), Inches(5.25), Inches(6.2), Inches(0.38),
    size=Pt(11), bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# 快速子項截圖（右）
rect(sl, Inches(6.9), Inches(1.2), Inches(6.2), Inches(4.2),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-quick-items.png", Inches(6.95), Inches(1.25), w=Inches(6.1))
txt(sl, "快速子項管理（48 項 / 8 類）",
    Inches(6.9), Inches(5.25), Inches(6.2), Inches(0.38),
    size=Pt(11), bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# 底部重點
pts = [
    "套用模板 → 30 秒建好完整 WBS",
    "拖曳快速子項 → 新增節點不需手動輸入",
    "科別可另存自訂模板，保留最佳實踐",
    "快速子項支援全域 + 科別兩種層級",
]
for i, p in enumerate(pts):
    x = Inches(0.38 + (i % 2) * 6.52)
    y = Inches(5.75 + (i // 2) * 0.5)
    rect(sl, x, y + Inches(0.1), Inches(0.06), Inches(0.26), fill=BLUE)
    txt(sl, p, x + Inches(0.14), y,
        Inches(6.2), Inches(0.44),
        size=Pt(12), color=BODY)

# ═══════════════════════════════════════════════════════════
# 7. 角色與權限
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "角色與權限", "四層組織架構，部門隔離，精細存取控制")
pnum(sl, 7)

hdrs  = ["角色",              "建立專案", "編輯 WBS",  "管理成員",  "跨科查閱"]
cws   = [Inches(3.4), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
cxs   = [Inches(0.45)]
for w in cws[:-1]: cxs.append(cxs[-1] + w + Inches(0.06))

rows_data = [
    ("DIRECTOR  部長",       "✅", "限自建",  "限自建",  "✅ 唯讀"),
    ("SECTION_CHIEF  科長",  "✅", "✅ 本科", "✅ 本科", "❌"),
    ("PROJECT_LEADER",       "✅", "✅ 加入", "✅ 自建", "❌"),
    ("PROJECT_MEMBER  成員", "❌", "✅ 加入", "❌",      "❌"),
]
for j, (h, cx, cw) in enumerate(zip(hdrs, cxs, cws)):
    rect(sl, cx, Inches(1.25), cw, Inches(0.52), fill=BLUE)
    txt(sl, h, cx, Inches(1.25), cw, Inches(0.52),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(rows_data):
    y = Inches(1.79 + i * 0.75)
    bg = BG if i % 2 == 0 else BG2
    for j, (val, cx, cw) in enumerate(zip(row, cxs, cws)):
        rect(sl, cx, y, cw, Inches(0.7), fill=bg, line_col=BORDER, lw=Pt(0.5))
        if j == 0:
            clr = DARK
        elif "✅" in val:
            clr = GREEN
        elif "❌" in val:
            clr = RED
        else:
            clr = AMBER
        txt(sl, val, cx + Inches(0.12), y,
            cw - Inches(0.12), Inches(0.7),
            size=Pt(12), bold=(j == 0),
            color=clr, align=(PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER))

# 底部說明框
rect(sl, Inches(0.45), Inches(5.0), Inches(12.5), Inches(1.35),
     fill=LBLUE, line_col=BLUE, lw=Pt(0.75))
txt(sl, "🔒  部門隔離保護",
    Inches(0.65), Inches(5.1), Inches(12), Inches(0.4),
    size=Pt(13), bold=True, color=BLUE)
txt(sl,
    "科長與 Leader 只能看到本科專案；部長可唯讀查閱下屬科，但不得編輯（自建除外）。\n"
    "歸檔後的專案對全員強制唯讀，系統以 canWriteProject() 作為唯一寫入判斷入口，確保一致性。",
    Inches(0.65), Inches(5.53), Inches(12), Inches(0.72),
    size=Pt(11.5), color=BODY)

# ═══════════════════════════════════════════════════════════
# 8. 與競品比較
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "與 GitHub 同類系統比較",
       "lgaleazzi/wbs（Spring Boot）  ·  Promess02/GanttApp（Node.js）  vs  本系統")
pnum(sl, 8)

c_rows = [
    ("WBS 節點管理",      "✅ 任意層樹狀",      "✅ 含甘特圖",    "✅ 固定 3 層"),
    ("即時多人協作",      "❌",                "❌",              "✅ WebSocket"),
    ("角色 / 部門權限",   "⚠️  基本登入",      "❌",             "✅ 4 層 + 隔離"),
    ("系統模板",          "❌",                "❌",              "✅ 4 套 + 自訂"),
    ("快速子項拖曳",      "❌",                "❌",              "✅ 48 項 / 8 類"),
    ("IT 上線標準流程",   "❌",                "❌",              "✅ SIT→UAT→PROD"),
    ("資料庫持久化",      "⚠️  H2（重啟清除）", "✅ SQLite",      "✅ PostgreSQL"),
    ("歸檔 + 匯出",       "❌",                "❌",              "✅ CSV / XLSX"),
]
c_hdrs  = ["功能",        "lgaleazzi/wbs",    "GanttApp",       "Spring-WbsScaff"]
c_cols  = [Inches(3.4), Inches(2.65), Inches(2.65), Inches(4.2)]
c_xs    = [Inches(0.38)]
for w in c_cols[:-1]: c_xs.append(c_xs[-1] + w + Inches(0.05))

h_bgs = [DARK, DARK, DARK, BLUE]
for j, (h, cx, cw, hbg) in enumerate(zip(c_hdrs, c_xs, c_cols, h_bgs)):
    rect(sl, cx, Inches(1.22), cw, Inches(0.5), fill=hbg)
    txt(sl, h, cx, Inches(1.22), cw, Inches(0.5),
        size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, row in enumerate(c_rows):
    y  = Inches(1.74 + i * 0.57)
    bg = BG if i % 2 == 0 else BG2
    for j, (val, cx, cw) in enumerate(zip(row, c_xs, c_cols)):
        fc, lc = check(val[:2] if val else "")
        cell_bg = LGREEN if (j==3 and "✅" in val) else (LRED if (j!=0 and "❌"==val.strip()) else (LAMBER if (j!=0 and "⚠️" in val) else bg))
        rect(sl, cx, y, cw, Inches(0.52), fill=cell_bg, line_col=BORDER, lw=Pt(0.4))
        clr = (DARK if j==0 else (GREEN if "✅" in val else (RED if "❌"==val.strip() else AMBER)))
        if j == 3: clr = BLUE if "✅" in val else clr
        txt(sl, val, cx + Inches(0.1), y,
            cw - Inches(0.14), Inches(0.52),
            size=Pt(11), bold=(j==3 and "✅" in val),
            color=clr, align=(PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER))

rect(sl, Inches(0.38), Inches(6.5), Inches(12.6), Inches(0.55),
     fill=LBLUE, line_col=BLUE, lw=Pt(0.75))
txt(sl,
    "同樣以 Spring Boot 為基底，Spring-WbsScaff 額外提供即時協作、精細權限與 IT 標準模板，專為企業 IT 部門設計。",
    Inches(0.55), Inches(6.52), Inches(12.3), Inches(0.48),
    size=Pt(11.5), color=BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 9. 功能驗證截圖
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "功能驗證", "登入 / 專案列表 / WBS 編輯器實際截圖")
pnum(sl, 9)

# 左：登入
rect(sl, Inches(0.38), Inches(1.2), Inches(4.4), Inches(5.7),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-login.png", Inches(0.43), Inches(1.25), w=Inches(4.3))
txt(sl, "✅  登入頁面",
    Inches(0.38), Inches(6.72), Inches(4.4), Inches(0.35),
    size=Pt(11), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# 右上：WBS
rect(sl, Inches(5.0), Inches(1.2), Inches(8.1), Inches(3.6),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-wbs-editor.png", Inches(5.05), Inches(1.25), w=Inches(8.0))
txt(sl, "✅  WBS 編輯器（功能新增模板，33 節點）",
    Inches(5.0), Inches(4.62), Inches(8.1), Inches(0.35),
    size=Pt(11), bold=True, color=GREEN)

# 右下：專案列表
rect(sl, Inches(5.0), Inches(5.1), Inches(8.1), Inches(1.85),
     fill=CARD, line_col=BORDER)
img(sl, "ppt-project-list.png", Inches(5.05), Inches(5.15), w=Inches(8.0))
txt(sl, "✅  專案列表（依角色過濾可見範圍）",
    Inches(5.0), Inches(6.78), Inches(8.1), Inches(0.35),
    size=Pt(11), bold=True, color=GREEN)

# ═══════════════════════════════════════════════════════════
# 10. 技術架構
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "技術架構")
pnum(sl, 10)

layers = [
    (BLUE,  "瀏覽器層",  "Thymeleaf SSR  +  Vue 3 CDN（離線靜態，無建置工具）"),
    (RGBColor(0x05,0x96,0x69), "通訊層", "REST API（X-CSRF-TOKEN）  +  SockJS / STOMP"),
    (RGBColor(0x7C,0x3A,0xED), "應用層", "Spring Boot 3.4  ·  Spring MVC  ·  Spring WebSocket  ·  Spring Security 6"),
    (RGBColor(0xD9,0x77,0x06), "資料層", "PostgreSQL 16  ·  Spring Session JDBC（HttpOnly Cookie，無 JWT）"),
]
for i, (col, lbl, desc) in enumerate(layers):
    y = Inches(1.3 + i * 1.35)
    rect(sl, Inches(0.38), y, Inches(12.6), Inches(1.1), fill=CARD, line_col=BORDER)
    rect(sl, Inches(0.38), y, Inches(0.07), Inches(1.1), fill=col)
    rect(sl, Inches(0.55), y + Inches(0.3), Inches(1.9), Inches(0.5), fill=col)
    txt(sl, lbl,
        Inches(0.55), y + Inches(0.3), Inches(1.9), Inches(0.5),
        size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc,
        Inches(2.6), y + Inches(0.3), Inches(10.2), Inches(0.5),
        size=Pt(13), color=DARK)
    if i < 3:
        txt(sl, "↓",
            Inches(6.6), y + Inches(1.1), Inches(0.4), Inches(0.22),
            size=Pt(11), color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 11. Roadmap：從原型到正式產品
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "從原型走向正式產品", "已完成的基礎，以及接下來的方向")
pnum(sl, 11)

phases = [
    ("Phase 1", "單頁原型", "已完成",
     ["WBS 節點編輯（CRUD）",
      "單一使用者操作",
      "基本的樹狀結構呈現",
      "解決 Excel 手工排版痛點"],
     GREEN, LGREEN),
    ("Phase 2", "協作平台", "現在",
     ["多人即時協作（WebSocket）",
      "角色權限 + 部門隔離",
      "4 套系統模板",
      "48 個快速子項 / 歸檔 / 匯出"],
     BLUE, LBLUE),
    ("Phase 3", "正式產品", "未來方向",
     ["行動裝置響應式介面",
      "通知系統（Email / Slack）",
      "甘特圖與里程碑視圖",
      "API 開放 + 第三方系統整合"],
     RGBColor(0x7C,0x3A,0xED), RGBColor(0xED,0xE9,0xFE)),
]
rect(sl, Inches(0.38), Inches(3.4), Inches(12.6), Pt(2), fill=BORDER)

for i, (phase, title, status, pts, col, lcol) in enumerate(phases):
    x = Inches(0.38 + i * 4.38)
    rect(sl, x + Inches(1.9), Inches(3.28), Inches(0.4), Inches(0.28), fill=col)

    rect(sl, x, Inches(1.2), Inches(4.15), Inches(5.9),
         fill=CARD if i < 2 else BG2, line_col=col, lw=Pt(1.5))
    rect(sl, x, Inches(1.2), Inches(4.15), Inches(0.08), fill=col)

    tag(sl, phase, x + Inches(0.18), Inches(1.35),
        Inches(1.3), Inches(0.36), bg=lcol, fg=col)
    txt(sl, title,
        x + Inches(1.62), Inches(1.33), Inches(2.4), Inches(0.42),
        size=Pt(16), bold=True, color=DARK)

    # 狀態
    s_bg = LGREEN if i==0 else (LBLUE if i==1 else RGBColor(0xED,0xE9,0xFE))
    s_fg = GREEN  if i==0 else (BLUE  if i==1 else RGBColor(0x7C,0x3A,0xED))
    tag(sl, status, x + Inches(0.18), Inches(1.85),
        Inches(2.0), Inches(0.34), bg=s_bg, fg=s_fg)

    rect(sl, x + Inches(0.18), Inches(2.33), Inches(3.8), Pt(1), fill=BORDER)

    for j, p in enumerate(pts):
        rect(sl, x + Inches(0.18), Inches(2.5 + j * 0.78),
             Inches(0.06), Inches(0.36), fill=col)
        txt(sl, p,
            x + Inches(0.35), Inches(2.48 + j * 0.78),
            Inches(3.65), Inches(0.72),
            size=Pt(12), color=BODY)

# ═══════════════════════════════════════════════════════════
# 12. 數字總結
# ═══════════════════════════════════════════════════════════
sl = new_slide(BG2)
header(sl, "一個平台，解決所有 WBS 管理需求")
pnum(sl, 12)

metrics = [
    ("4",    "套系統模板",   "功能新增、EOS 下線\n專案開發、系統開發",    BLUE),
    ("48",   "個快速子項",   "分 8 類，涵蓋所有\nIT 標準工作項目",        GREEN),
    ("3",    "層 WBS 結構",  "L1 大項 → L2 工作包\n→ L3 任務（含備註）",  RGBColor(0x7C,0x3A,0xED)),
    ("即時", "多人協作",     "WebSocket 廣播\n秒級同步，零刷新",          RGBColor(0xD9,0x77,0x06)),
]
for i, (num, unit, desc, col) in enumerate(metrics):
    x = Inches(0.38 + i * 3.25)
    rect(sl, x, Inches(1.2), Inches(3.05), Inches(5.7),
         fill=BG, line_col=col, lw=Pt(1.5))
    rect(sl, x, Inches(1.2), Inches(3.05), Inches(0.08), fill=col)
    txt(sl, num,
        x, Inches(1.55), Inches(3.05), Inches(1.5),
        size=Pt(62), bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(sl, unit,
        x, Inches(3.15), Inches(3.05), Inches(0.55),
        size=Pt(17), bold=True, color=DARK, align=PP_ALIGN.CENTER)
    rect(sl, x + Inches(0.3), Inches(3.8), Inches(2.4), Pt(1), fill=BORDER)
    txt(sl, desc,
        x + Inches(0.12), Inches(3.95), Inches(2.82), Inches(1.5),
        size=Pt(12), color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# 13. 快速啟動
# ═══════════════════════════════════════════════════════════
sl = new_slide()
header(sl, "立即開始使用")
pnum(sl, 13)

# 左側：指令
rect(sl, Inches(0.38), Inches(1.2), Inches(7.4), Inches(5.9),
     fill=DARK, line_col=BORDER)
rect(sl, Inches(0.38), Inches(1.2), Inches(7.4), Inches(0.42), fill=RGBColor(0x1E,0x40,0x6E))
txt(sl, "  Terminal",
    Inches(0.45), Inches(1.22), Inches(3), Inches(0.38),
    size=Pt(11), color=MUTED)

cmds = [
    ("# 複製環境設定",               False),
    ("cp .env.example .env",         True),
    ("",                             False),
    ("# 啟動 PostgreSQL",            False),
    ("docker compose up -d",         True),
    ("",                             False),
    ("# 編譯並啟動",                 False),
    ("mvn clean install -DskipTests",True),
    ("mvn spring-boot:run",          True),
    ("",                             False),
    ("# 開啟瀏覽器",                 False),
    ("open http://localhost:8080",   True),
]
y_off = Inches(1.72)
for line, is_cmd in cmds:
    clr = WHITE if is_cmd else GREEN
    txt(sl, line or " ",
        Inches(0.6), y_off, Inches(7.0), Inches(0.36),
        size=Pt(12), color=clr)
    y_off += Inches(0.38 if is_cmd else 0.3)

# 右側：帳號
txt(sl, "測試帳號  ( 密碼皆為  test1234 )",
    Inches(8.1), Inches(1.2), Inches(5.0), Inches(0.45),
    size=Pt(13), bold=True, color=DARK)

accounts = [
    ("director@company.com",  "DIRECTOR  部長",      AMBER),
    ("chief@infotech.com",    "SECTION_CHIEF  科長", BLUE),
    ("leader@infotech.com",   "PROJECT_LEADER",      GREEN),
    ("member1@infotech.com",  "PROJECT_MEMBER",      MUTED),
]
for i, (email, role, col) in enumerate(accounts):
    y = Inches(1.8 + i * 0.98)
    rect(sl, Inches(8.1), y, Inches(5.0), Inches(0.88),
         fill=CARD, line_col=col, lw=Pt(1.2))
    rect(sl, Inches(8.1), y, Inches(0.06), Inches(0.88), fill=col)
    txt(sl, email, Inches(8.25), y + Inches(0.08),
        Inches(4.7), Inches(0.38),
        size=Pt(12), color=DARK)
    txt(sl, role, Inches(8.25), y + Inches(0.48),
        Inches(4.7), Inches(0.34),
        size=Pt(11), color=col)

rect(sl, Inches(8.1), Inches(5.85), Inches(5.0), Inches(1.25),
     fill=LBLUE, line_col=BLUE, lw=Pt(1))
txt(sl, "🚀  從單頁原型出發",
    Inches(8.3), Inches(5.98), Inches(4.6), Inches(0.42),
    size=Pt(14), bold=True, color=BLUE)
txt(sl, "每一行程式都是為了讓 IT 團隊\n不再被 WBS 的雜亂拖慢腳步。",
    Inches(8.3), Inches(6.42), Inches(4.6), Inches(0.6),
    size=Pt(11.5), color=BODY)

prs.save(OUT)
print(f"✅  已輸出：{OUT}  （共 {TOTAL} 張）")
